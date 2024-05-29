# -*- coding: UTF-8 -*-
import os
import sys
from pptx import Presentation
import translators as ts

def is_number(s):
    try:
        float(s)
        return True
    except ValueError:
        return False

def replace_text_in_pptx(input_file, output_file, translations=None, src_lang='auto', dest_lang='en'):
    # 设置工作目录
    if getattr(sys, 'frozen', False):
        bundle_dir = sys._MEIPASS
    else:
        bundle_dir = os.path.dirname(os.path.abspath(__file__))
    os.chdir(bundle_dir)

    # 加载原始PPTX文件
    prs = Presentation(input_file)
    
    # 翻译文本（如果还没翻译的话）
    if not translations:
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    text = ''.join(run.text for paragraph in shape.text_frame.paragraphs for run in paragraph.runs)
                    texts.append(text)
                elif shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            text = ''.join(para.text for para in cell.text_frame.paragraphs)
                            texts.append(text)

        
        translations = {}  # 初始化一个空字典来存储翻译结果
        for text in texts:
            if not is_number(text):
                translated_text = ts.translate_text(text, to_language="zh") 
                translations[text] = translated_text  # 将原文和翻译后的文本存入字典
                print(text,translated_text)
            else:
                translations[text] = text
            
    
    # 遍历幻灯片和形状，替换文本
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                text = ''.join(run.text for paragraph in shape.text_frame.paragraphs for run in paragraph.runs)
                if text in translations:
                    # 清除原有文本
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.text = ''
                    # 添加翻译后的文本
                    shape.text_frame.text = translations[text]
            elif shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        text = ''.join(para.text for para in cell.text_frame.paragraphs)
                        if text in translations:
                            # 清除原有文本
                            for para in cell.text_frame.paragraphs:
                                para.clear()
                            # 添加翻译后的文本
                            cell.text_frame.text = translations[text]
    
    # 保存到新文件
    prs.save(output_file)

# 使用示例
# translations_dict = {...}  # 例如：{'原文1': '翻译后的内容1', '原文2': '翻译后的内容2'}
# replace_text_in_pptx('original.pptx', 'translated.pptx', translations_dict)

# 或者直接调用进行翻译和替换
replace_text_in_pptx('DV発煙不具合報告20240417.pptx', 'translated.pptx', None)
